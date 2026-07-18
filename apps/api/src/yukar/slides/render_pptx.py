"""Deck → editable .pptx, in-process via python-pptx (no subprocess).

Geometry contract: canvas px map 1:1 to EMU at 96 dpi (9525 EMU/px), so the
numbers in the definition are exactly what PowerPoint stores.  Line spacing
is written as an *exact* point height (``size * line_height``) rather than a
multiple — PowerPoint's "multiple" spacing depends on font metrics, which the
HTML preview could not reproduce; an exact height renders identically in
both.

python-pptx has no bullet API and no East-Asian font API; both are done with
small lxml edits on the underlying XML (well-known recipes).  Missing images
render as a labelled placeholder box so the deck stays openable; the caller
reports the warning.
"""

from __future__ import annotations

import io
import posixpath
from collections.abc import Mapping
from dataclasses import dataclass
from typing import Any

from lxml import etree  # ty: ignore[unresolved-import] — lxml ships no stubs
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from yukar.slides.schema import (
    Deck,
    ImageElement,
    LineElement,
    Paragraph,
    ShapeElement,
    Slide,
    TableElement,
    TextElement,
)

EMU_PER_PX = 9525  # 96 dpi
_BULLET_INDENT_EMU = 24 * EMU_PER_PX  # 24 px per level, shared with the CSS preview
_BLANK_LAYOUT_INDEX = 6  # "Blank" in python-pptx's default template

# Shared with render_html so placeholder boxes look the same in both outputs.
PLACEHOLDER_FILL = "#EEEEEE"
PLACEHOLDER_LINE = "#BBBBBB"
PLACEHOLDER_TEXT_COLOR = "#888888"
PLACEHOLDER_FONT_PT = 10.0

TABLE_ZEBRA_FILL = "#F2F2F2"
TABLE_BODY_FILL = "#FFFFFF"

ROUNDED_CORNER_FRACTION = 0.12  # of the shorter side, mirrored as CSS border-radius

_ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
_ANCHOR = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
_SHAPES = {
    "rect": MSO_SHAPE.RECTANGLE,
    "rounded": MSO_SHAPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_SHAPE.OVAL,
}


@dataclass(frozen=True, slots=True)
class LoadedImage:
    """Image bytes plus natural pixel size (needed for contain/cover math)."""

    data: bytes
    width: int
    height: int


def image_key(path: str) -> str:
    """Canonical dict key for an image path.

    Spelling variants of the same file (``x.png`` / ``./x.png`` /
    ``a//x.png``) must collapse to one key, so the loader reads each file
    once and both renderers find it under the same entry.
    """
    return posixpath.normpath(path)


def _emu(px: float) -> Emu:
    return Emu(round(px * EMU_PER_PX))


def _rgb(color: str) -> RGBColor:
    return RGBColor.from_string(color.lstrip("#"))


def _set_run_style(
    run: Any,
    *,
    font_name: str | None,
    size_pt: float,
    bold: bool,
    italic: bool,
    color: str,
) -> None:
    """Style one run, including the East-Asian typeface python-pptx omits.

    ``font.name`` writes only ``<a:latin>``; CJK glyphs then fall back to the
    theme font and diverge from the preview.  ``<a:ea>``/``<a:cs>`` follow
    ``<a:latin>`` in the schema sequence, so appending is order-safe here
    (no hlink/sym elements are ever present on our runs).
    """
    font = run.font
    font.size = Pt(size_pt)
    font.bold = bold
    font.italic = italic
    font.color.rgb = _rgb(color)
    if font_name:
        font.name = font_name
        r_pr = run._r.get_or_add_rPr()  # noqa: SLF001 — no public EA-font API
        for tag in ("a:ea", "a:cs"):
            el = r_pr.find(qn(tag))
            if el is None:
                el = etree.SubElement(r_pr, qn(tag))
            el.set("typeface", font_name)


def _apply_bullet(paragraph: Any, level: int) -> None:
    """Real ``<a:buChar>`` bullets (python-pptx has no bullet API).

    Indent geometry matches the preview CSS: hanging indent of 24 px, list
    body at ``24 * (level + 1)`` px.  Appending is order-safe because the
    paragraph carries at most lnSpc/spcBef before this runs, and bu* elements
    sort after those in ``CT_TextParagraphProperties``.
    """
    paragraph.level = level
    p_pr = paragraph._p.get_or_add_pPr()  # noqa: SLF001 — no public bullet API
    p_pr.set("marL", str(_BULLET_INDENT_EMU * (level + 1)))
    p_pr.set("indent", str(-_BULLET_INDENT_EMU))
    bu_font = etree.SubElement(p_pr, qn("a:buFont"))
    bu_font.set("typeface", "Arial")
    bu_char = etree.SubElement(p_pr, qn("a:buChar"))
    bu_char.set("char", "•")


def _add_paragraph(text_frame: Any, index: int) -> Any:
    return text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()


def _add_text_runs(
    p: Any,
    text: str,
    *,
    font_name: str | None,
    size_pt: float,
    bold: bool,
    italic: bool,
    color: str,
) -> None:
    """Write *text* as styled runs, turning ``\\n`` into real ``<a:br/>`` breaks.

    ``run.text`` would store a literal newline inside ``<a:t>``, which
    PowerPoint renders as whitespace, not a break; the HTML preview uses
    ``white-space: pre-wrap``, so both sides need ``\\n`` to be a hard break.
    """
    for j, segment in enumerate(text.split("\n")):
        if j:
            etree.SubElement(p._p, qn("a:br"))  # noqa: SLF001 — no public <a:br> API
        run = p.add_run()
        run.text = segment
        _set_run_style(
            run, font_name=font_name, size_pt=size_pt, bold=bold, italic=italic, color=color
        )


def _render_paragraph(
    p: Any, para: Paragraph, deck: Deck, default_align: str
) -> None:
    size = para.size if para.size is not None else deck.font_size
    p.alignment = _ALIGN[para.align or default_align]
    p.line_spacing = Pt(size * para.line_height)
    if para.space_before:
        p.space_before = Pt(para.space_before)
    if para.bullet:
        _apply_bullet(p, para.level)
    _add_text_runs(
        p,
        para.text,
        font_name=deck.font,
        size_pt=size,
        bold=para.bold,
        italic=para.italic,
        color=para.color or deck.text_color,
    )


def _render_text(slide_obj: Any, el: TextElement, deck: Deck) -> None:
    box = slide_obj.shapes.add_textbox(_emu(el.x), _emu(el.y), _emu(el.w), _emu(el.h))
    tf = box.text_frame
    tf.word_wrap = True
    # add_textbox defaults to spAutoFit (shape grows to fit text) — the box
    # must keep its declared size so overflow behaves like the preview.
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = _ANCHOR[el.valign]
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    for i, para in enumerate(el.paragraphs):
        _render_paragraph(_add_paragraph(tf, i), para, deck, el.align)


def _render_missing_image(slide_obj: Any, el: ImageElement, deck: Deck) -> None:
    """Grey labelled box where the image would be — deck stays presentable."""
    sp = slide_obj.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, _emu(el.x), _emu(el.y), _emu(el.w), _emu(el.h)
    )
    sp.shadow.inherit = False
    sp.fill.solid()
    sp.fill.fore_color.rgb = _rgb(PLACEHOLDER_FILL)
    sp.line.color.rgb = _rgb(PLACEHOLDER_LINE)
    sp.line.width = Pt(0.75)
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _add_text_runs(
        p,
        f"image not found: {el.path}",
        font_name=deck.font,
        size_pt=PLACEHOLDER_FONT_PT,
        bold=False,
        italic=False,
        color=PLACEHOLDER_TEXT_COLOR,
    )


def _render_image(
    slide_obj: Any, el: ImageElement, deck: Deck, images: Mapping[str, LoadedImage]
) -> None:
    loaded = images.get(image_key(el.path))
    if loaded is None:
        _render_missing_image(slide_obj, el, deck)
        return
    stream = io.BytesIO(loaded.data)
    if el.fit == "stretch" or loaded.width <= 0 or loaded.height <= 0:
        slide_obj.shapes.add_picture(stream, _emu(el.x), _emu(el.y), _emu(el.w), _emu(el.h))
        return
    if el.fit == "contain":
        scale = min(el.w / loaded.width, el.h / loaded.height)
        w2, h2 = loaded.width * scale, loaded.height * scale
        slide_obj.shapes.add_picture(
            stream,
            _emu(el.x + (el.w - w2) / 2),
            _emu(el.y + (el.h - h2) / 2),
            _emu(w2),
            _emu(h2),
        )
        return
    # cover: fill the box and crop the overhang symmetrically (srcRect).
    scale = max(el.w / loaded.width, el.h / loaded.height)
    visible_w = el.w / (loaded.width * scale)
    visible_h = el.h / (loaded.height * scale)
    pic = slide_obj.shapes.add_picture(
        stream, _emu(el.x), _emu(el.y), _emu(el.w), _emu(el.h)
    )
    pic.crop_left = pic.crop_right = max(0.0, (1 - visible_w) / 2)
    pic.crop_top = pic.crop_bottom = max(0.0, (1 - visible_h) / 2)


def _render_shape(slide_obj: Any, el: ShapeElement) -> None:
    sp = slide_obj.shapes.add_shape(
        _SHAPES[el.shape], _emu(el.x), _emu(el.y), _emu(el.w), _emu(el.h)
    )
    sp.shadow.inherit = False
    if el.shape == "rounded":
        sp.adjustments[0] = ROUNDED_CORNER_FRACTION
    if el.fill is not None:
        sp.fill.solid()
        sp.fill.fore_color.rgb = _rgb(el.fill)
    else:
        sp.fill.background()
    if el.line_color is not None:
        sp.line.color.rgb = _rgb(el.line_color)
        sp.line.width = Pt(el.line_width * 0.75)  # px → pt
    else:
        sp.line.fill.background()


def _render_line(slide_obj: Any, el: LineElement) -> None:
    conn = slide_obj.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, _emu(el.x1), _emu(el.y1), _emu(el.x2), _emu(el.y2)
    )
    conn.shadow.inherit = False
    conn.line.color.rgb = _rgb(el.color)
    conn.line.width = Pt(el.width * 0.75)


def _cell_fill(el: TableElement, row_index: int) -> str:
    if el.header and row_index == 0:
        return el.header_fill
    body_index = row_index - (1 if el.header else 0)
    return TABLE_ZEBRA_FILL if el.zebra and body_index % 2 == 1 else TABLE_BODY_FILL


def _render_table(slide_obj: Any, el: TableElement, deck: Deck) -> None:
    n_rows, n_cols = len(el.rows), len(el.rows[0])
    frame = slide_obj.shapes.add_table(
        n_rows, n_cols, _emu(el.x), _emu(el.y), _emu(el.w), _emu(el.h)
    )
    table = frame.table
    table.first_row = el.header
    table.horz_banding = False  # zebra is written explicitly for preview parity

    weights = el.col_widths or [1.0] * n_cols
    total = sum(weights)
    consumed = 0
    for c in range(n_cols):
        width = (
            _emu(el.w) - consumed
            if c == n_cols - 1
            else _emu(el.w * weights[c] / total)
        )
        # Extreme ratios can round a column (or the last column's remainder)
        # to zero or below — clamp so the written gridCol is always valid.
        width = max(width, 1)
        table.columns[c].width = Emu(width)
        consumed += width
    for r in range(n_rows):
        table.rows[r].height = _emu(el.h / n_rows)

    size = el.font_size if el.font_size is not None else deck.font_size
    for r, row in enumerate(el.rows):
        is_header = el.header and r == 0
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(_cell_fill(el, r))
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Emu(8 * EMU_PER_PX)
            cell.margin_top = cell.margin_bottom = Emu(4 * EMU_PER_PX)
            _add_text_runs(
                cell.text_frame.paragraphs[0],
                text,
                font_name=deck.font,
                size_pt=size,
                bold=is_header,
                italic=False,
                color=el.header_color if is_header else deck.text_color,
            )


def _render_slide(
    slide_obj: Any, slide: Slide, deck: Deck, images: Mapping[str, LoadedImage]
) -> None:
    fill = slide_obj.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(slide.background or deck.background)
    for el in slide.elements:
        if isinstance(el, TextElement):
            _render_text(slide_obj, el, deck)
        elif isinstance(el, ImageElement):
            _render_image(slide_obj, el, deck, images)
        elif isinstance(el, ShapeElement):
            _render_shape(slide_obj, el)
        elif isinstance(el, LineElement):
            _render_line(slide_obj, el)
        elif isinstance(el, TableElement):
            _render_table(slide_obj, el, deck)
    if slide.notes:
        slide_obj.notes_slide.notes_text_frame.text = slide.notes


def render_pptx(deck: Deck, images: Mapping[str, LoadedImage]) -> bytes:
    """Render the validated deck to .pptx bytes (pure CPU, no I/O)."""
    prs = Presentation()
    width, height = deck.canvas
    prs.slide_width = _emu(width)
    prs.slide_height = _emu(height)
    blank = prs.slide_layouts[_BLANK_LAYOUT_INDEX]
    for slide in deck.slides:
        _render_slide(prs.slides.add_slide(blank), slide, deck, images)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
